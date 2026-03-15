from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OCEAN = RGBColor(0x0E, 0x74, 0x90)
OCEAN_DARK = RGBColor(0x0A, 0x5B, 0x72)
CORAL = RGBColor(0xF9, 0x73, 0x16)
SAND = RGBColor(0xFB, 0xF7, 0xF0)
CHARCOAL = RGBColor(0x1C, 0x19, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DRIFTWOOD = RGBColor(0xC2, 0x95, 0x6B)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        sld_fill = shape.fill._fill
        solid = sld_fill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solid is not None:
            srgb = solid.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                alpha_el = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha_el.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=CHARCOAL):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
        p.level = 0
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=OCEAN, font_size=11):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ''
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Calibri'
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = CHARCOAL
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table_shape


def slide_header(slide, title, subtitle=None):
    add_shape_bg(slide, 0, 0, W, Inches(1.4), OCEAN)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6), title, 32, WHITE, True, font_name='Calibri')
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.4), subtitle, 16, RGBColor(0xBB, 0xDD, 0xEE), font_name='Calibri')


# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, OCEAN_DARK)

add_shape_bg(slide1, 0, Inches(0.5), W, Inches(0.08), CORAL)

add_text_box(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             'ILOT COCO BEACH', 54, WHITE, True, PP_ALIGN.CENTER, 'Calibri')
add_text_box(slide1, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             'Projet Web Full-Stack', 32, CORAL, False, PP_ALIGN.CENTER, 'Calibri')
add_text_box(slide1, Inches(1), Inches(3.7), Inches(11), Inches(0.6),
             'Proposition Technique & Commerciale', 20, RGBColor(0xBB, 0xDD, 0xEE), False, PP_ALIGN.CENTER)

add_shape_bg(slide1, Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.04), DRIFTWOOD)

add_text_box(slide1, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             'Mars 2026', 18, RGBColor(0x88, 0xBB, 0xCC), False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: PROJECT OVERVIEW
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, SAND)
slide_header(slide2, 'Vue d\'ensemble du projet', 'Ce que nous construisons')

items = [
    'Site web public vitrine avec video de couverture et galerie photo',
    'Systeme de reservation avec integration WhatsApp',
    'Menu digital avec deux versions : Standard et Extra',
    'Systeme de vouchers et ventes flash avec compte a rebours',
    'Inscription utilisateurs (email/telephone + Google/Facebook)',
    'Systeme d\'avis clients avec moderation admin',
    'Backoffice complet pour gerer tout le contenu',
    'Integration reseaux sociaux : WhatsApp, Instagram, Messenger',
    'Site multilingue : Francais, Arabe, Anglais',
    'Cron jobs pour activation/desactivation automatique des promos',
]
add_bullet_list(slide2, Inches(0.8), Inches(1.7), Inches(7), Inches(5.5), items, 16, CHARCOAL)

# Right side summary box
box = add_shape_bg(slide2, Inches(8.5), Inches(1.8), Inches(4.2), Inches(4.5), WHITE)
add_text_box(slide2, Inches(8.8), Inches(1.9), Inches(3.8), Inches(0.5), 'Chiffres cles', 20, OCEAN, True)
stats = [
    '40 User Stories',
    '10 Tables de base de donnees',
    '3 Langues supportees',
    '~200 heures de developpement',
    '10-12 semaines de livraison',
    '27 fonctionnalites client',
    '13 fonctionnalites backoffice',
    '2 cron jobs automatiques',
]
add_bullet_list(slide2, Inches(8.8), Inches(2.5), Inches(3.8), Inches(3.5), stats, 14, CHARCOAL)


# ============================================================
# SLIDE 3: FEATURE MAP - PUBLIC SITE
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, SAND)
slide_header(slide3, 'Fonctionnalites - Site Public', 'Ce que vos clients voient')

data = [
    ['#', 'Fonctionnalite', 'Description'],
    ['1', 'Video de couverture (Hero)', 'Video plein ecran en arriere-plan avec identite de marque'],
    ['2', 'A propos / Decouvrir', 'Presentation de l\'activite et atouts uniques'],
    ['3', 'Menu / Plats', 'Deux versions : Standard et Extra, avec prix et images'],
    ['4', 'Galerie photo', 'Grille masonry geree depuis le backoffice'],
    ['5', 'Emplacements', 'Cabanes, pied dans l\'eau, plage, balade pirate en mer'],
    ['6', 'Vouchers / Promos', 'Codes promo et reductions affiches sur le site'],
    ['7', 'Vente Flash', 'Ventes limitees avec compte a rebours et reduction'],
    ['8', 'Avis clients', 'Temoignages approuves par l\'admin'],
    ['9', 'Reseaux sociaux', 'Boutons WhatsApp, Instagram, Messenger'],
    ['10', 'Localisation', 'Carte Google Maps, telephone, itineraire'],
    ['11', 'Multilingue', 'Francais, Arabe (RTL), Anglais'],
]
add_table(slide3, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.2), 12, 3, data,
          col_widths=[0.5, 3.0, 8.8], font_size=12)


# ============================================================
# SLIDE 4: FEATURE MAP - BACKOFFICE
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, SAND)
slide_header(slide4, 'Fonctionnalites - Backoffice Admin', 'Ce que vous gerez vous-meme')

data = [
    ['#', 'Fonctionnalite', 'Description'],
    ['1', 'Tableau de bord', 'Statistiques : visites, reservations, promos actives'],
    ['2', 'Gestion du menu', 'Ajouter/modifier/supprimer plats, prix Standard/Extra, images'],
    ['3', 'Gestion des images', 'Upload, reorganiser, supprimer pour toutes les sections'],
    ['4', 'Gestion des emplacements', 'Descriptions, images, prix, capacite'],
    ['5', 'Gestion des vouchers', 'Creer des codes promo avec reduction, date d\'expiration'],
    ['6', 'Gestion ventes flash', 'Programmer debut/fin, reduction %, activation auto'],
    ['7', 'Moderation des avis', 'Voir, approuver, rejeter ou masquer les avis clients'],
    ['8', 'Reservations recues', 'Consulter les demandes de reservation'],
    ['9', 'Gestion du contenu', 'Modifier textes, video hero, infos contact, mots-cles SEO'],
    ['10', 'Gestion utilisateurs', 'Voir les utilisateurs inscrits'],
]
add_table(slide4, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.8), 11, 3, data,
          col_widths=[0.5, 3.5, 8.3], font_size=12)

add_text_box(slide4, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
             'Tout est modifiable sans intervention technique : photos, prix, menus, promos, textes.',
             14, OCEAN, True)

# ============================================================
# SLIDE 5: ARCHITECTURE
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, SAND)
slide_header(slide5, 'Architecture Technique', 'Comment le systeme est construit')

# Frontend box
box1 = add_shape_bg(slide5, Inches(0.5), Inches(2.0), Inches(3.5), Inches(3.5), WHITE)
add_text_box(slide5, Inches(0.7), Inches(2.1), Inches(3.1), Inches(0.5), 'FRONTEND', 20, OCEAN, True, PP_ALIGN.CENTER)
items_fe = ['Vue.js 3 + Vite', 'TailwindCSS v4', 'Vue Router + Vue I18n', 'GSAP Animations', 'Lenis Smooth Scroll', 'Pinia (state)']
add_bullet_list(slide5, Inches(0.8), Inches(2.7), Inches(3.0), Inches(2.5), items_fe, 13)

# Arrow
add_text_box(slide5, Inches(4.1), Inches(3.3), Inches(0.8), Inches(0.5), '<-->', 24, CORAL, True, PP_ALIGN.CENTER)

# Backend box
box2 = add_shape_bg(slide5, Inches(4.9), Inches(2.0), Inches(3.5), Inches(3.5), WHITE)
add_text_box(slide5, Inches(5.1), Inches(2.1), Inches(3.1), Inches(0.5), 'BACKEND', 20, OCEAN, True, PP_ALIGN.CENTER)
items_be = ['Node.js (Fastify)', 'Prisma ORM', 'JWT Authentication', 'node-cron (jobs)', 'Multer (file upload)', 'REST API']
add_bullet_list(slide5, Inches(5.2), Inches(2.7), Inches(3.0), Inches(2.5), items_be, 13)

# Arrow
add_text_box(slide5, Inches(8.5), Inches(3.3), Inches(0.8), Inches(0.5), '<-->', 24, CORAL, True, PP_ALIGN.CENTER)

# DB box
box3 = add_shape_bg(slide5, Inches(9.3), Inches(2.0), Inches(3.5), Inches(3.5), WHITE)
add_text_box(slide5, Inches(9.5), Inches(2.1), Inches(3.1), Inches(0.5), 'BASE DE DONNEES', 20, OCEAN, True, PP_ALIGN.CENTER)
items_db = ['PostgreSQL', '10 tables', 'Users & Auth', 'Menus & Items', 'Vouchers & Flash Sales', 'Reviews & Reservations']
add_bullet_list(slide5, Inches(9.6), Inches(2.7), Inches(3.0), Inches(2.5), items_db, 13)

# Storage box bottom
box4 = add_shape_bg(slide5, Inches(2.5), Inches(5.8), Inches(3.5), Inches(1.2), WHITE)
add_text_box(slide5, Inches(2.7), Inches(5.9), Inches(3.1), Inches(0.4), 'STOCKAGE FICHIERS', 16, OCEAN, True, PP_ALIGN.CENTER)
add_text_box(slide5, Inches(2.7), Inches(6.3), Inches(3.1), Inches(0.5), 'Cloudflare R2 / S3\nImages, videos', 12, CHARCOAL, False, PP_ALIGN.CENTER)

# Hosting box bottom
box5 = add_shape_bg(slide5, Inches(7.3), Inches(5.8), Inches(3.5), Inches(1.2), WHITE)
add_text_box(slide5, Inches(7.5), Inches(5.9), Inches(3.1), Inches(0.4), 'HEBERGEMENT', 16, OCEAN, True, PP_ALIGN.CENTER)
add_text_box(slide5, Inches(7.5), Inches(6.3), Inches(3.1), Inches(0.5), 'Railway / Hetzner / OVH\nVercel (frontend)', 12, CHARCOAL, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6: BACKEND FRAMEWORK COMPARISON
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, SAND)
slide_header(slide6, 'Comparaison des Frameworks Backend', 'Quel langage choisir ?')

data = [
    ['Critere', 'Node.js', 'Django (Python)', '.NET (C#)', 'Laravel (PHP)'],
    ['Vitesse de dev', '★★★★★', '★★★★★', '★★★☆☆', '★★★★☆'],
    ['Cout hebergement', '★★★★★', '★★★★☆', '★★★☆☆', '★★★★★'],
    ['Meme langage frontend', '★★★★★', '★★☆☆☆', '★★☆☆☆', '★★☆☆☆'],
    ['Ecosysteme', '★★★★★', '★★★★★', '★★★★☆', '★★★★☆'],
    ['Cron Jobs', '★★★★★', '★★★★★', '★★★★☆', '★★★★★'],
    ['Migration OVH', '★★★★☆', '★★★★☆', '★★★☆☆', '★★★★★'],
    ['TOTAL', '34/35', '29/35', '22/35', '29/35'],
]
add_table(slide6, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.8), 8, 5, data,
          col_widths=[2.8, 2.3, 2.3, 2.3, 2.3], font_size=13)

# Recommendation box
rec_box = add_shape_bg(slide6, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.5), OCEAN_DARK)
add_text_box(slide6, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
             'RECOMMANDATION : Node.js (Fastify + Prisma + PostgreSQL)', 22, WHITE, True)
add_text_box(slide6, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7),
             'Un seul langage (JavaScript) pour tout le projet = un seul developpeur peut tout gerer. '
             'Hebergement le moins cher. Developpement le plus rapide pour les operations CRUD.', 14, RGBColor(0xBB, 0xDD, 0xEE))


# ============================================================
# SLIDE 7: CLOUD PROVIDER COMPARISON
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, SAND)
slide_header(slide7, 'Comparaison Hebergement Cloud', 'Options et couts mensuels')

add_text_box(slide7, Inches(0.5), Inches(1.5), Inches(6), Inches(0.4), 'Option A : Plateformes gerees (plus simple)', 16, OCEAN, True)
data_a = [
    ['Fournisseur', 'Backend', 'Base de donnees', 'Cout/mois'],
    ['Railway', 'Inclus ($5 credit)', 'PostgreSQL inclus', '~$5-10'],
    ['Render', 'Starter $7/mois', 'PostgreSQL $7/mois', '~$7-14'],
    ['Vercel + Supabase', 'Serverless (gratuit)', 'Supabase gratuit', '$0 (tier gratuit)'],
]
add_table(slide7, Inches(0.5), Inches(1.95), Inches(6), Inches(2.0), 4, 4, data_a,
          col_widths=[1.5, 1.5, 1.5, 1.2], font_size=11)

add_text_box(slide7, Inches(7), Inches(1.5), Inches(6), Inches(0.4), 'Option B : VPS (moins cher long terme)', 16, OCEAN, True)
data_b = [
    ['Fournisseur', 'Spec', 'Prix/mois'],
    ['Hetzner CX22', '2 vCPU, 4GB RAM, 40GB', '3.79 EUR'],
    ['OVH VPS', '1 vCPU, 2GB RAM, 20GB', '3.50 EUR'],
    ['DigitalOcean', '1 vCPU, 1GB RAM, 25GB', '$6.00'],
    ['Contabo VPS', '4 vCPU, 8GB RAM, 50GB', '4.99 EUR'],
]
add_table(slide7, Inches(7), Inches(1.95), Inches(5.8), Inches(2.3), 5, 3, data_b,
          col_widths=[1.8, 2.2, 1.3], font_size=11)

# Recommended
rec = add_shape_bg(slide7, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), WHITE)
add_text_box(slide7, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.5), 'Plan recommande', 20, OCEAN, True)

add_text_box(slide7, Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.3), 'Phase 1 - Lancement', 15, CORAL, True)
add_text_box(slide7, Inches(0.8), Inches(5.6), Inches(5.5), Inches(1.2),
             'Frontend : Vercel (gratuit)\nBackend : Railway ($5/mois)\nBase de donnees : incluse\nStockage : Cloudflare R2 (gratuit)\n\nTotal : ~$5-7/mois', 12, CHARCOAL)

add_text_box(slide7, Inches(7), Inches(5.3), Inches(5.5), Inches(0.3), 'Phase 2 - Migration OVH', 15, CORAL, True)
add_text_box(slide7, Inches(7), Inches(5.6), Inches(5.5), Inches(1.2),
             'Tout sur Hetzner CX22 ou OVH VPS\nPostgreSQL auto-heberge\nNode.js + Nginx + Docker\nCloudflare CDN (gratuit)\n\nTotal : ~4-7 EUR/mois', 12, CHARCOAL)


# ============================================================
# SLIDE 8: USER STORIES / BACKLOG
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, SAND)
slide_header(slide8, 'Backlog - User Stories', 'Decoupage du travail en epics')

data = [
    ['Epic', 'Nb Stories', 'Heures', 'Sprints'],
    ['Site public (hero, menu, galerie, emplacements)', '12', '50h', '2'],
    ['Authentification utilisateurs', '4', '17h', '1'],
    ['Systeme de reservation', '3', '11h', '0.5'],
    ['Avis & feedback', '2', '6h', '0.5'],
    ['Backoffice - Gestion contenu', '6', '30h', '1.5'],
    ['Backoffice - Vouchers & Ventes flash', '4', '18h', '1'],
    ['Backoffice - Moderation', '4', '11h', '0.5'],
    ['Infrastructure backend', '5', '26h', '1.5'],
    ['Tests & QA', '-', '20h', '1'],
    ['Deploiement & lancement', '-', '10h', '0.5'],
    ['TOTAL', '40', '~199h', '10'],
]
add_table(slide8, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.2), 12, 4, data,
          col_widths=[5.5, 1.5, 1.5, 1.5], font_size=13)


# ============================================================
# SLIDE 9: TIMELINE
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, SAND)
slide_header(slide9, 'Planning de developpement', 'Timeline estimee : 10-12 semaines')

phases = [
    ('Semaines 1-2', 'Infrastructure + Site Public', 'Setup backend, DB, API de base\nHero, About, Galerie, Menu frontend', OCEAN),
    ('Semaines 3-4', 'Auth + Contenu dynamique', 'Inscription/connexion utilisateurs\nEmplacements, avis, reseaux sociaux', RGBColor(0x06, 0x96, 0x68)),
    ('Semaines 5-6', 'Backoffice - Contenu', 'Admin : menus, images, emplacements\nGestion du contenu du site', DRIFTWOOD),
    ('Semaines 7-8', 'Vouchers & Ventes Flash', 'Systeme de vouchers + ventes flash\nCron jobs + comptes a rebours', CORAL),
    ('Semaines 9-10', 'Tests & Deploiement', 'QA, corrections, responsive\nMise en production + formation', RGBColor(0x7C, 0x3A, 0xED)),
]

for i, (period, title, desc, color) in enumerate(phases):
    y = Inches(1.8) + Inches(i * 1.1)
    # Colored bar
    add_shape_bg(slide9, Inches(0.5), y, Inches(0.15), Inches(0.9), color)
    # Period
    add_text_box(slide9, Inches(0.9), y, Inches(2.2), Inches(0.4), period, 16, color, True)
    # Title
    add_text_box(slide9, Inches(0.9), y + Inches(0.35), Inches(2.2), Inches(0.4), title, 13, CHARCOAL, True)
    # Description
    add_text_box(slide9, Inches(3.5), y, Inches(9), Inches(0.9), desc, 13, CHARCOAL)


# ============================================================
# SLIDE 10: COST ESTIMATION
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, SAND)
slide_header(slide10, 'Estimation des couts', 'Developpement + Hebergement')

# Dev cost
add_text_box(slide10, Inches(0.5), Inches(1.5), Inches(6), Inches(0.4), 'A. Cout de developpement', 18, OCEAN, True)
data_dev = [
    ['Approche', 'Taux', 'Heures', 'Total'],
    ['Freelance (marche tunisien)', '30-50 TND/h', '200h', '6,000-10,000 TND'],
    ['Freelance (international)', '$25-50/h', '200h', '$5,000-10,000'],
    ['Petite agence', '$40-80/h', '200h', '$8,000-16,000'],
]
add_table(slide10, Inches(0.5), Inches(2.0), Inches(6), Inches(2.0), 4, 4, data_dev,
          col_widths=[2.0, 1.3, 0.8, 2.2], font_size=12)

# Hosting cost
add_text_box(slide10, Inches(7), Inches(1.5), Inches(6), Inches(0.4), 'B. Cout hebergement mensuel', 18, OCEAN, True)
data_host = [
    ['Phase', 'Composants', 'Cout/mois'],
    ['Lancement', 'Railway + Vercel + CDN', '$5-10'],
    ['Production', 'VPS Hetzner/OVH + Cloudflare', '4-7 EUR'],
    ['Domaine', '.com / .tn', '~$1/mois'],
    ['TOTAL', '', '$6-15/mois'],
]
add_table(slide10, Inches(7), Inches(2.0), Inches(5.8), Inches(2.3), 5, 3, data_host,
          col_widths=[1.5, 2.5, 1.3], font_size=12)

# Suggested pricing
prop_box = add_shape_bg(slide10, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.3), OCEAN_DARK)
add_text_box(slide10, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.5),
             'Proposition tarifaire suggeree au client', 22, WHITE, True)

add_text_box(slide10, Inches(0.8), Inches(5.5), Inches(3.8), Inches(1.3),
             'Option 1 : Forfait\n\n10,000 TND\n+ 200 TND/mois maintenance', 16, WHITE)
add_text_box(slide10, Inches(4.8), Inches(5.5), Inches(3.8), Inches(1.3),
             'Option 2 : Phasage\n\nPhase 1 (MVP) : 6,000 TND\nPhase 2 (complet) : 4,000 TND', 16, WHITE)
add_text_box(slide10, Inches(8.8), Inches(5.5), Inches(3.8), Inches(1.3),
             'Inclus :\n\nCode source + 3 mois support\nFormation backoffice\nDocumentation', 16, RGBColor(0xBB, 0xDD, 0xEE))


# ============================================================
# SLIDE 11: DELIVERABLES
# ============================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11, SAND)
slide_header(slide11, 'Livrables & Autonomie client', 'Ce que le client recoit et gere seul')

# What client gets
add_text_box(slide11, Inches(0.5), Inches(1.6), Inches(6), Inches(0.4), 'Le client recoit :', 18, OCEAN, True)
gets = [
    'Site web responsive (desktop, tablette, mobile)',
    'Backoffice complet pour gestion autonome',
    'Systeme d\'inscription et d\'avis clients',
    'Moteur de vouchers et ventes flash automatise',
    'Support multilingue (FR/AR/EN)',
    'Integration WhatsApp, Instagram, Messenger',
    'SEO optimise avec mots-cles configurables',
    'Propriete du code source',
    'Documentation d\'utilisation du backoffice',
    '3 mois de support post-lancement',
]
add_bullet_list(slide11, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5), gets, 13, CHARCOAL)

# What client manages
add_text_box(slide11, Inches(7), Inches(1.6), Inches(6), Inches(0.4), 'Le client gere seul :', 18, GREEN, True)
manages = [
    'Toutes les photos et videos',
    'Les plats du menu et les prix',
    'Les vouchers et promotions',
    'La programmation des ventes flash',
    'La moderation des avis clients',
    'Les informations de contact et SEO',
]
add_bullet_list(slide11, Inches(7.1), Inches(2.1), Inches(5.5), Inches(3), manages, 13, CHARCOAL)

add_text_box(slide11, Inches(7), Inches(4.5), Inches(6), Inches(0.4), 'Necessite un developpeur :', 18, RED, True)
dev_needs = [
    'Nouvelles fonctionnalites',
    'Changements majeurs de design',
    'Migration de serveur',
    'Ajout de nouvelles langues',
]
add_bullet_list(slide11, Inches(7.1), Inches(5.0), Inches(5.5), Inches(2), dev_needs, 13, CHARCOAL)


# ============================================================
# SLIDE 12: NEXT STEPS
# ============================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12, OCEAN_DARK)

add_text_box(slide12, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             'Prochaines etapes', 40, WHITE, True, PP_ALIGN.CENTER)

add_shape_bg(slide12, Inches(3), Inches(1.7), Inches(7), Inches(0.04), CORAL)

steps = [
    ('01', 'Validation du perimetre et du budget par le client'),
    ('02', 'Creation des comptes cloud (Railway/Hetzner + base de donnees)'),
    ('03', 'Sprint 1-2 : Infrastructure backend + site public frontend'),
    ('04', 'Sprint 3-4 : Authentification + Menu + Galerie + Emplacements'),
    ('05', 'Sprint 5-6 : Backoffice (contenu, menus, images)'),
    ('06', 'Sprint 7-8 : Vouchers, ventes flash, avis, cron jobs'),
    ('07', 'Sprint 9-10 : Tests, QA, deploiement, lancement'),
    ('08', 'Post-lancement : Suivi, corrections, optimisation SEO'),
]

for i, (num, text) in enumerate(steps):
    y = Inches(2.2) + Inches(i * 0.6)
    # Number circle
    circle = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CORAL
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide12, Inches(2.2), y + Inches(0.05), Inches(9), Inches(0.4), text, 17, WHITE)


# ============================================================
# SLIDE 13: THANK YOU
# ============================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide13, OCEAN_DARK)

add_text_box(slide13, Inches(1), Inches(2.0), Inches(11), Inches(1),
             'ILOT COCO BEACH', 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide13, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
             'Merci pour votre confiance', 28, CORAL, False, PP_ALIGN.CENTER)
add_shape_bg(slide13, Inches(5), Inches(4.2), Inches(3.3), Inches(0.04), DRIFTWOOD)
add_text_box(slide13, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
             'Des questions ?', 22, RGBColor(0xBB, 0xDD, 0xEE), False, PP_ALIGN.CENTER)


output_path = r'C:\Personal\ilot-cocobeach\Ilot_CocoBeach_Presentation.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
