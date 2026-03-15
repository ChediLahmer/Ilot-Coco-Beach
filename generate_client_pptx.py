from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

OCEAN = RGBColor(0x0E, 0x74, 0x90)
OCEAN_DARK = RGBColor(0x08, 0x4D, 0x63)
OCEAN_DEEPER = RGBColor(0x05, 0x33, 0x44)
CORAL = RGBColor(0xF9, 0x73, 0x16)
CORAL_LIGHT = RGBColor(0xFD, 0xBA, 0x74)
SAND = RGBColor(0xFB, 0xF7, 0xF0)
SAND_DARK = RGBColor(0xF0, 0xE6, 0xD3)
CHARCOAL = RGBColor(0x1C, 0x19, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DRIFTWOOD = RGBColor(0xC2, 0x95, 0x6B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
LIGHT_OCEAN = RGBColor(0xCC, 0xFA, 0xFA)
MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)

nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def gradient_bg(slide, color1, color2, angle=270):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0
    try:
        grad_xml = fill._fill._element.find(f'{{{nsmap["a"]}}}gradFill')
        if grad_xml is not None:
            lin = grad_xml.find(f'{{{nsmap["a"]}}}lin')
            if lin is None:
                lin = etree.SubElement(grad_xml, f'{{{nsmap["a"]}}}lin')
            lin.set('ang', str(angle * 60000))
            lin.set('scaled', '1')
    except Exception:
        pass


def solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT, font='Calibri', italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.font.italic = italic
    p.alignment = align
    p.space_after = Pt(0)
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=14, color=CHARCOAL, line_spacing=1.3, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(size * 0.4)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def card(slide, left, top, width, height, title, body, icon_text=None, title_color=OCEAN, accent_color=OCEAN):
    r = add_rect(slide, left, top, width, height, WHITE, radius=0.06)
    # Accent bar top
    add_rect(slide, left + Inches(0.02), top + Inches(0.02), width - Inches(0.04), Inches(0.06), accent_color)
    # Icon circle
    if icon_text:
        c = add_circle(slide, left + Inches(0.3), top + Inches(0.35), Inches(0.55), accent_color)
        tf = c.text_frame
        tf.paragraphs[0].text = icon_text
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_left = left + Inches(1.05)
        title_width = width - Inches(1.35)
    else:
        title_left = left + Inches(0.3)
        title_width = width - Inches(0.6)
    add_text(slide, title_left, top + Inches(0.35), title_width, Inches(0.4), title, 16, title_color, True)
    add_text(slide, left + Inches(0.3), top + Inches(0.9), width - Inches(0.6), height - Inches(1.1), body, 12, MUTED)
    return r


def section_title(slide, title, subtitle=None, light=False):
    tc = WHITE if light else OCEAN
    sc = RGBColor(0xAA, 0xCC, 0xDD) if light else MUTED
    # Decorative line
    line_color = CORAL
    add_rect(slide, Inches(0.8), Inches(1.0), Inches(0.6), Inches(0.05), line_color)
    add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7), title, 36, tc, True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5), subtitle, 16, sc, italic=True)


def page_number(slide, num, total, light=False):
    c = RGBColor(0x88, 0xAA, 0xBB) if light else MUTED
    add_text(slide, Inches(12), Inches(7.0), Inches(1), Inches(0.4), f'{num}/{total}', 10, c, align=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12

# ================================================================
# SLIDE 1: TITLE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
gradient_bg(s, OCEAN_DARK, OCEAN_DEEPER, 315)

# Decorative circles (top right, bottom left)
from lxml import etree as ET

def set_shape_alpha(shape, alpha_val):
    try:
        fill_elem = shape.fill._fill._element
        solid = fill_elem.find(f'{{{nsmap["a"]}}}solidFill')
        if solid is not None:
            srgb = solid.find(f'{{{nsmap["a"]}}}srgbClr')
            if srgb is not None:
                alpha_el = ET.SubElement(srgb, f'{{{nsmap["a"]}}}alpha')
                alpha_el.set('val', str(alpha_val))
    except Exception:
        pass

c1 = add_circle(s, Inches(10.5), Inches(-1.5), Inches(4), OCEAN)
set_shape_alpha(c1, 15000)

c2 = add_circle(s, Inches(-2), Inches(5), Inches(5), OCEAN)
set_shape_alpha(c2, 10000)

# Decorative line
add_rect(s, Inches(4.5), Inches(2.5), Inches(4.3), Inches(0.05), CORAL)

add_text(s, Inches(1), Inches(1.2), Inches(11.3), Inches(1.2), 'ILOT COCO BEACH', 56, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.8), Inches(11.3), Inches(0.8), 'Votre presence digitale, a la hauteur de votre plage', 22, CORAL_LIGHT, False, PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6), 'Proposition de projet web', 20, RGBColor(0x88, 0xBB, 0xCC), False, PP_ALIGN.CENTER)

# Bottom bar
add_rect(s, 0, Inches(6.6), W, Inches(0.9), OCEAN_DEEPER)
add_text(s, Inches(1), Inches(6.75), Inches(5), Inches(0.4), 'Mars 2026', 14, RGBColor(0x66, 0x99, 0xAA))
add_text(s, Inches(7), Inches(6.75), Inches(5.3), Inches(0.4), 'Confidentiel', 14, RGBColor(0x66, 0x99, 0xAA), align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 2: THE VISION
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
section_title(s, 'Notre vision pour vous', 'Un site web qui reflete la beaute de votre plage')
page_number(s, 2, TOTAL_SLIDES)

# Three vision pillars
pillars = [
    ('Attirer', 'Un design qui capte l\'attention des visiteurs des le premier regard. Video immersive, galerie photo, ambiance paradise.', '01', OCEAN),
    ('Convertir', 'Systeme de reservation simple, ventes flash avec compte a rebours, vouchers promotionnels pour transformer les visiteurs en clients.', '02', CORAL),
    ('Fideliser', 'Espace client avec avis, programme de reductions, et communication directe via WhatsApp, Instagram et Messenger.', '03', GREEN),
]

for i, (title, body, num, color) in enumerate(pillars):
    x = Inches(0.8) + Inches(i * 4.1)
    card(s, x, Inches(2.2), Inches(3.7), Inches(3.8), title, body, num, CHARCOAL, color)

# Bottom note
add_rect(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7), OCEAN_DARK, radius=0.08)
add_text(s, Inches(1.2), Inches(6.4), Inches(11), Inches(0.5),
         'Tout le contenu du site sera geerable par vous-meme, sans aucune competence technique requise.',
         14, WHITE, False, PP_ALIGN.CENTER)

# ================================================================
# SLIDE 3: WHAT YOUR CLIENTS SEE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
section_title(s, 'Ce que vos clients decouvrent', 'L\'experience de navigation sur votre site')
page_number(s, 3, TOTAL_SLIDES)

features = [
    ('Video immersive', 'Une video plein ecran en page d\'accueil qui plonge le visiteur dans l\'ambiance de votre plage.', OCEAN),
    ('Menu digital', 'Vos plats presentes avec photos, en deux formules : Standard et Extra. Prix et disponibilite a jour.', CORAL),
    ('Galerie photo', 'Vos plus belles images dans une galerie interactive et elegante, mise a jour depuis le backoffice.', OCEAN),
    ('Emplacements', 'Cabanes, pied dans l\'eau, plage, balade pirate en mer. Chaque espace avec sa description et tarif.', DRIFTWOOD),
    ('Ventes flash', 'Offres limitees avec compte a rebours visible. Creez l\'urgence et boostez vos reservations.', RGBColor(0xEF, 0x44, 0x44)),
    ('Vouchers', 'Codes promotionnels affiches sur le site. Vos clients les appliquent lors de la reservation.', GREEN),
    ('Avis clients', 'Les temoignages de vos clients satisfaits, que vous choisissez d\'afficher ou non.', OCEAN),
    ('Contact direct', 'Boutons WhatsApp, Instagram et Messenger pour une communication instantanee.', RGBColor(0x25, 0xD3, 0x66)),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(col * 3.15)
    y = Inches(2.0) + Inches(row * 2.7)
    r = add_rect(s, x, y, Inches(2.95), Inches(2.4), WHITE, radius=0.06)
    add_rect(s, x, y, Inches(0.07), Inches(2.4), color)
    add_text(s, x + Inches(0.3), y + Inches(0.25), Inches(2.4), Inches(0.35), title, 15, CHARCOAL, True)
    add_text(s, x + Inches(0.3), y + Inches(0.7), Inches(2.4), Inches(1.5), desc, 11, MUTED)

# ================================================================
# SLIDE 4: BACKOFFICE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
gradient_bg(s, OCEAN_DARK, OCEAN_DEEPER, 315)
section_title(s, 'Votre tableau de bord', 'Gerez tout votre contenu en autonomie', light=True)
page_number(s, 4, TOTAL_SLIDES, light=True)

admin_features = [
    ('Menu & Prix', 'Ajoutez, modifiez ou supprimez vos plats.\nGerez les prix Standard et Extra.\nActivez/desactivez des plats.'),
    ('Photos & Videos', 'Ajoutez de nouvelles photos.\nReorganisez votre galerie.\nChangez la video d\'accueil.'),
    ('Vouchers', 'Creez des codes promo avec\nreduction en % ou montant fixe.\nDefinissez les dates de validite.'),
    ('Ventes Flash', 'Programmez des ventes limitees.\nLe systeme les active et desactive\nautomatiquement a l\'heure prevue.'),
    ('Avis clients', 'Consultez tous les avis recus.\nApprouvez ceux a afficher.\nMasquez ou supprimez les autres.'),
    ('Reservations', 'Recevez les demandes en temps reel.\nConsultez l\'historique.\nRepondez via WhatsApp.'),
]

for i, (title, desc) in enumerate(admin_features):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(2.0) + Inches(row * 2.7)
    r = add_rect(s, x, y, Inches(3.9), Inches(2.4), RGBColor(0x0C, 0x5E, 0x76), radius=0.06)
    # accent top
    add_rect(s, x + Inches(0.02), y + Inches(0.02), Inches(3.86), Inches(0.05), CORAL)
    add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.35), title, 17, WHITE, True)
    add_text(s, x + Inches(0.3), y + Inches(0.8), Inches(3.3), Inches(1.5), desc, 12, RGBColor(0xAA, 0xCC, 0xDD))

# ================================================================
# SLIDE 5: HOW IT WORKS (USER JOURNEY)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
section_title(s, 'Le parcours de votre client', 'De la decouverte a la reservation')
page_number(s, 5, TOTAL_SLIDES)

steps = [
    ('1', 'Decouvre\nvotre site', 'Via Google, reseaux\nsociaux ou bouche\na oreille', OCEAN),
    ('2', 'Explore\nle contenu', 'Video, galerie,\nmenu, emplacements\net avis clients', RGBColor(0x06, 0x96, 0x68)),
    ('3', 'Voit une\npromotion', 'Vente flash ou\nvoucher qui cree\nl\'envie de venir', CORAL),
    ('4', 'Reserve\nsa place', 'Formulaire simple\nenvoye directement\nvia WhatsApp', RGBColor(0x7C, 0x3A, 0xED)),
    ('5', 'Partage\nson avis', 'Apres sa visite,\nlaisse un temoignage\nsur le site', DRIFTWOOD),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.4) + Inches(i * 2.55)
    # Card
    add_rect(s, x, Inches(2.5), Inches(2.25), Inches(4.0), WHITE, radius=0.06)
    # Number circle
    c = add_circle(s, x + Inches(0.75), Inches(2.2), Inches(0.7), color)
    tf = c.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    add_text(s, x + Inches(0.15), Inches(3.1), Inches(1.95), Inches(0.7), title, 16, CHARCOAL, True, PP_ALIGN.CENTER)
    # Desc
    add_text(s, x + Inches(0.15), Inches(4.0), Inches(1.95), Inches(1.5), desc, 12, MUTED, align=PP_ALIGN.CENTER)
    # Arrow between steps
    if i < len(steps) - 1:
        add_text(s, x + Inches(2.25), Inches(4.0), Inches(0.3), Inches(0.5), '>', 28, CORAL, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 6: MULTILINGUE + SOCIAL
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, WHITE)
section_title(s, 'Touchez plus de clients', 'Multilingue + Reseaux sociaux')
page_number(s, 6, TOTAL_SLIDES)

# Language cards
langs = [
    ('Francais', 'Langue principale du site.\nTout le contenu, menus et\ninterface en francais.', '🇫🇷'),
    ('Arabe', 'Interface completement\nadaptee de droite a gauche\n(RTL) pour vos clients arabes.', '🇹🇳'),
    ('Anglais', 'Pour les touristes\ninternationaux qui visitent\nla region.', '🇬🇧'),
]

for i, (title, desc, flag) in enumerate(langs):
    x = Inches(0.8) + Inches(i * 2.8)
    add_rect(s, x, Inches(2.0), Inches(2.5), Inches(2.2), SAND, radius=0.06)
    add_text(s, x + Inches(0.2), Inches(2.15), Inches(2.1), Inches(0.4), f'{flag}  {title}', 16, CHARCOAL, True, PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.65), Inches(2.1), Inches(1.3), desc, 12, MUTED, align=PP_ALIGN.CENTER)

# Social cards
socials = [
    ('WhatsApp', 'Reservation directe et\ncommunication en temps reel\navec vos clients.', RGBColor(0x25, 0xD3, 0x66)),
    ('Instagram', 'Redirigez vers votre page\nInstagram pour partager\nvos moments.', RGBColor(0xE1, 0x30, 0x6C)),
    ('Messenger', 'Discussion directe via\nFacebook Messenger pour\nles questions rapides.', RGBColor(0x00, 0x84, 0xFF)),
]

for i, (title, desc, color) in enumerate(socials):
    x = Inches(0.8) + Inches(i * 2.8)
    add_rect(s, x, Inches(4.6), Inches(2.5), Inches(2.2), WHITE, radius=0.06)
    add_rect(s, x, Inches(4.6), Inches(2.5), Inches(0.06), color)
    add_text(s, x + Inches(0.2), Inches(4.8), Inches(2.1), Inches(0.4), title, 16, color, True, PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(5.3), Inches(2.1), Inches(1.3), desc, 12, MUTED, align=PP_ALIGN.CENTER)

# Right side image placeholder
add_rect(s, Inches(9.3), Inches(2.0), Inches(3.5), Inches(4.8), SAND, radius=0.08)
add_text(s, Inches(9.5), Inches(3.8), Inches(3.1), Inches(1.0),
         'Votre site s\'adapte\nautomatiquement a\nchaque langue et\nchaque ecran', 16, OCEAN, True, PP_ALIGN.CENTER)

# ================================================================
# SLIDE 7: INVESTMENT - DEVELOPMENT
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
section_title(s, 'Votre investissement', 'Un projet cle en main')
page_number(s, 7, TOTAL_SLIDES)

# Two pricing options
# Option 1
add_rect(s, Inches(1), Inches(2.2), Inches(5), Inches(4.5), WHITE, radius=0.06)
add_rect(s, Inches(1), Inches(2.2), Inches(5), Inches(0.8), OCEAN, radius=0.06)
add_rect(s, Inches(1), Inches(2.7), Inches(5), Inches(0.4), OCEAN)
add_text(s, Inches(1.3), Inches(2.3), Inches(4.4), Inches(0.6), 'FORFAIT COMPLET', 22, WHITE, True, PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(3.3), Inches(4), Inches(0.8), '10,000 TND', 40, OCEAN, True, PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.1), Inches(4), Inches(0.4), '+ 200 TND / mois maintenance', 14, CORAL, True, PP_ALIGN.CENTER)

inclus1 = [
    'Site web complet avec toutes les fonctionnalites',
    'Backoffice de gestion autonome',
    'Multilingue (FR / AR / EN)',
    'Hebergement 1ere annee inclus',
    'Formation d\'utilisation',
    '3 mois de support post-lancement',
    'Code source : votre propriete',
]
for i, item in enumerate(inclus1):
    y = Inches(4.55) + Inches(i * 0.28)
    add_text(s, Inches(1.6), y, Inches(3.8), Inches(0.25), f'  {item}', 11, CHARCOAL)
    add_circle(s, Inches(1.45), y + Inches(0.04), Inches(0.12), GREEN)

# Option 2
add_rect(s, Inches(7), Inches(2.2), Inches(5), Inches(4.5), WHITE, radius=0.06)
add_rect(s, Inches(7), Inches(2.2), Inches(5), Inches(0.8), CORAL, radius=0.06)
add_rect(s, Inches(7), Inches(2.7), Inches(5), Inches(0.4), CORAL)
add_text(s, Inches(7.3), Inches(2.3), Inches(4.4), Inches(0.6), 'EN DEUX PHASES', 22, WHITE, True, PP_ALIGN.CENTER)

add_text(s, Inches(7.5), Inches(3.3), Inches(4), Inches(0.4), 'Phase 1 (MVP)', 16, CHARCOAL, True, PP_ALIGN.CENTER)
add_text(s, Inches(7.5), Inches(3.7), Inches(4), Inches(0.5), '6,000 TND', 32, CORAL, True, PP_ALIGN.CENTER)
add_text(s, Inches(7.5), Inches(4.2), Inches(4), Inches(0.3), 'Site public + Menu + Galerie + Reservation', 11, MUTED, align=PP_ALIGN.CENTER)

add_text(s, Inches(7.5), Inches(4.7), Inches(4), Inches(0.4), 'Phase 2 (Complet)', 16, CHARCOAL, True, PP_ALIGN.CENTER)
add_text(s, Inches(7.5), Inches(5.1), Inches(4), Inches(0.5), '4,000 TND', 32, CORAL, True, PP_ALIGN.CENTER)
add_text(s, Inches(7.5), Inches(5.6), Inches(4), Inches(0.3), 'Backoffice + Vouchers + Flash + Users + Avis', 11, MUTED, align=PP_ALIGN.CENTER)

add_text(s, Inches(7.5), Inches(6.1), Inches(4), Inches(0.35), '+ 200 TND / mois maintenance', 14, CORAL, True, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 8: MONTHLY COSTS
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, WHITE)
section_title(s, 'Couts de fonctionnement', 'Ce que coute l\'hebergement chaque mois')
page_number(s, 8, TOTAL_SLIDES)

items_cost = [
    ('Hebergement serveur', 'Le serveur qui fait tourner votre site et votre backoffice, accessible 24h/24.', '4-7 EUR', OCEAN),
    ('Nom de domaine', 'Votre adresse web personnalisee (ex: ilotcocobeach.com)', '~1 EUR', DRIFTWOOD),
    ('Certificat SSL', 'Le cadenas de securite sur votre site. Gratuit avec Let\'s Encrypt.', 'Gratuit', GREEN),
    ('CDN & Securite', 'Acceleration du chargement et protection contre les attaques. Cloudflare gratuit.', 'Gratuit', GREEN),
]

for i, (title, desc, cost, color) in enumerate(items_cost):
    y = Inches(2.0) + Inches(i * 1.2)
    add_rect(s, Inches(0.8), y, Inches(8.5), Inches(1.0), SAND, radius=0.06)
    add_rect(s, Inches(0.8), y, Inches(0.08), Inches(1.0), color)
    add_text(s, Inches(1.2), y + Inches(0.1), Inches(4), Inches(0.35), title, 15, CHARCOAL, True)
    add_text(s, Inches(1.2), y + Inches(0.45), Inches(6.5), Inches(0.4), desc, 11, MUTED)
    add_text(s, Inches(7.5), y + Inches(0.2), Inches(1.5), Inches(0.5), cost, 18, color, True, PP_ALIGN.RIGHT)

# Total box
add_rect(s, Inches(0.8), Inches(6.2), Inches(8.5), Inches(0.8), OCEAN_DARK, radius=0.08)
add_text(s, Inches(1.2), Inches(6.3), Inches(4), Inches(0.5), 'TOTAL MENSUEL', 18, WHITE, True)
add_text(s, Inches(6), Inches(6.3), Inches(3), Inches(0.5), '5 - 8 EUR / mois', 22, CORAL_LIGHT, True, PP_ALIGN.RIGHT)

# Comparison note
add_rect(s, Inches(10), Inches(2.0), Inches(2.8), Inches(5.0), SAND, radius=0.08)
add_text(s, Inches(10.2), Inches(2.2), Inches(2.4), Inches(0.5), 'A titre de\ncomparaison', 14, OCEAN, True, PP_ALIGN.CENTER)
comparisons = [
    'Un site Wix Pro\ncoute $17/mois',
    'Un site Squarespace\ncoute $27/mois',
    'Une page Facebook\nsponsored coute\n$50+/mois',
    'Votre site :\n~6 EUR/mois\nSANS abonnement',
]
y_comp = Inches(3.0)
for comp in comparisons:
    add_text(s, Inches(10.3), y_comp, Inches(2.2), Inches(1.0), comp, 11, MUTED, align=PP_ALIGN.CENTER)
    y_comp += Inches(1.1)


# ================================================================
# SLIDE 9: TIMELINE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
section_title(s, 'Calendrier de livraison', 'Votre site pret en 10 semaines')
page_number(s, 9, TOTAL_SLIDES)

phases = [
    ('Sem 1-2', 'Fondations', 'Mise en place du site,\nvideo d\'accueil, galerie', OCEAN),
    ('Sem 3-4', 'Contenu', 'Menu digital, emplacements,\navis clients, reservation', RGBColor(0x06, 0x96, 0x68)),
    ('Sem 5-6', 'Backoffice', 'Panneau d\'administration\npour gerer tout le contenu', DRIFTWOOD),
    ('Sem 7-8', 'Promotions', 'Vouchers, ventes flash,\ncron jobs automatiques', CORAL),
    ('Sem 9-10', 'Lancement', 'Tests, corrections,\nmise en ligne + formation', RGBColor(0x7C, 0x3A, 0xED)),
]

# Horizontal timeline
line_y = Inches(3.6)
add_rect(s, Inches(1.2), line_y, Inches(10.9), Inches(0.04), DRIFTWOOD)

for i, (period, title, desc, color) in enumerate(phases):
    x = Inches(1.2) + Inches(i * 2.7)
    # Circle on line
    c = add_circle(s, x + Inches(0.7), line_y - Inches(0.2), Inches(0.45), color)
    tf = c.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Period label above
    add_text(s, x, Inches(2.4), Inches(1.9), Inches(0.3), period, 12, color, True, PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.7), Inches(1.9), Inches(0.4), title, 16, CHARCOAL, True, PP_ALIGN.CENTER)
    # Description below line
    add_text(s, x, Inches(4.2), Inches(1.9), Inches(1.2), desc, 12, MUTED, align=PP_ALIGN.CENTER)

# Delivery note
add_rect(s, Inches(2), Inches(5.8), Inches(9.3), Inches(1.0), WHITE, radius=0.08)
add_text(s, Inches(2.3), Inches(5.9), Inches(8.7), Inches(0.35), 'Vous recevez des mises a jour a chaque etape.', 15, OCEAN, True, PP_ALIGN.CENTER)
add_text(s, Inches(2.3), Inches(6.3), Inches(8.7), Inches(0.35), 'Vos retours sont integres au fur et a mesure pour un resultat qui vous ressemble.', 13, MUTED, align=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 10: WHAT YOU GET
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
gradient_bg(s, OCEAN_DARK, OCEAN_DEEPER, 315)
section_title(s, 'Ce que vous recevez', 'Tout ce qui est inclus dans le projet', light=True)
page_number(s, 10, TOTAL_SLIDES, light=True)

deliverables = [
    'Site web responsive : parfait sur ordinateur, tablette et telephone',
    'Backoffice complet : gerez menus, photos, prix, promos en autonomie',
    'Systeme de reservation via WhatsApp : zero cout par message',
    'Ventes flash automatisees : programmez, le systeme fait le reste',
    'Multilingue : Francais, Arabe et Anglais',
    'Avis clients : vos clients temoignent, vous moderez',
    'SEO optimise : mots-cles configurables pour etre trouve sur Google',
    'Code source : votre propriete, vous n\'etes dependant de personne',
    'Formation : nous vous montrons comment tout gerer',
    '3 mois de support : nous restons disponibles apres le lancement',
]

for i, item in enumerate(deliverables):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(2.0) + Inches(row * 1.0)
    # check circle
    c = add_circle(s, x, y + Inches(0.08), Inches(0.3), GREEN)
    tf = c.text_frame
    tf.paragraphs[0].text = '✓'
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text(s, x + Inches(0.5), y, Inches(5.3), Inches(0.5), item, 14, WHITE)


# ================================================================
# SLIDE 11: NEXT STEPS
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
section_title(s, 'Prochaines etapes', 'Comment on avance ensemble')
page_number(s, 11, TOTAL_SLIDES)

next_steps = [
    ('Validation', 'Vous validez le perimetre du projet et le budget.', OCEAN),
    ('Lancement', 'Nous demarrons le developpement immediatement.', CORAL),
    ('Suivi', 'Vous recevez des demos regulieres a chaque etape.', RGBColor(0x06, 0x96, 0x68)),
    ('Livraison', 'Le site est mis en ligne et vous etes forme.', RGBColor(0x7C, 0x3A, 0xED)),
]

for i, (title, desc, color) in enumerate(next_steps):
    x = Inches(0.8) + Inches(i * 3.15)
    add_rect(s, x, Inches(2.5), Inches(2.85), Inches(3.0), WHITE, radius=0.06)
    c = add_circle(s, x + Inches(1.05), Inches(2.2), Inches(0.7), color)
    tf = c.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text(s, x + Inches(0.2), Inches(3.2), Inches(2.45), Inches(0.4), title, 18, CHARCOAL, True, PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(3.7), Inches(2.45), Inches(1.2), desc, 13, MUTED, align=PP_ALIGN.CENTER)

add_rect(s, Inches(2.5), Inches(6.0), Inches(8.3), Inches(0.9), OCEAN_DARK, radius=0.08)
add_text(s, Inches(2.8), Inches(6.15), Inches(7.7), Inches(0.5),
         'Pret a demarrer ? Contactez-nous pour fixer un rendez-vous de lancement.',
         16, WHITE, False, PP_ALIGN.CENTER)

# ================================================================
# SLIDE 12: THANK YOU
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
gradient_bg(s, OCEAN_DARK, OCEAN_DEEPER, 315)

c1 = add_circle(s, Inches(9.5), Inches(-2), Inches(5), OCEAN)
set_shape_alpha(c1, 12000)

add_rect(s, Inches(4.8), Inches(2.3), Inches(3.7), Inches(0.05), CORAL)

add_text(s, Inches(1), Inches(1.0), Inches(11.3), Inches(1), 'ILOT COCO BEACH', 50, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.6), Inches(11.3), Inches(0.7), 'Merci pour votre confiance', 26, CORAL_LIGHT, False, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.6), 'Votre ile, votre evasion.', 20, RGBColor(0x88, 0xBB, 0xCC), False, PP_ALIGN.CENTER, italic=True)

add_rect(s, Inches(3.5), Inches(5.0), Inches(6.3), Inches(1.8), RGBColor(0x0C, 0x5E, 0x76), radius=0.08)
add_text(s, Inches(3.8), Inches(5.2), Inches(5.7), Inches(0.4), 'Des questions ? Contactez-nous', 16, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(3.8), Inches(5.7), Inches(5.7), Inches(0.4), 'contact@ilotcocobeach.com', 14, CORAL_LIGHT, align=PP_ALIGN.CENTER)
add_text(s, Inches(3.8), Inches(6.1), Inches(5.7), Inches(0.4), '+216 99 123 456', 14, RGBColor(0xAA, 0xCC, 0xDD), align=PP_ALIGN.CENTER)


output = r'C:\Personal\ilot-cocobeach\Ilot_CocoBeach_CLIENT.pptx'
prs.save(output)
print(f'Client presentation saved: {output}')
print(f'Slides: {len(prs.slides)}')
