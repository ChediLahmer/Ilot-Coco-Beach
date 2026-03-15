from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OCEAN = RGBColor(0x0E, 0x74, 0x90)
OCEAN_DARK = RGBColor(0x08, 0x4D, 0x63)
OCEAN_DEEPER = RGBColor(0x05, 0x33, 0x44)
CORAL = RGBColor(0xF9, 0x73, 0x16)
CORAL_LIGHT = RGBColor(0xFD, 0xBA, 0x74)
SAND = RGBColor(0xFB, 0xF7, 0xF0)
CHARCOAL = RGBColor(0x1C, 0x19, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DRIFTWOOD = RGBColor(0xC2, 0x95, 0x6B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
MUTED = RGBColor(0x6B, 0x72, 0x80)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
CODE_BG = RGBColor(0x1E, 0x29, 0x3B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def circle(slide, l, t, sz, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, l, t, w, h, text, sz=18, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT, italic=False, font='Calibri'):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.font.italic = italic
    p.alignment = align
    p.space_after = Pt(0)
    return txBox


def multiline(slide, l, t, w, h, lines, sz=13, color=CHARCOAL):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(sz * 0.3)
    return txBox


def header(slide, title, sub=None, dark=False):
    tc = WHITE if dark else OCEAN
    sc = RGBColor(0xAA, 0xCC, 0xDD) if dark else MUTED
    rect(slide, Inches(0.8), Inches(1.0), Inches(0.6), Inches(0.05), CORAL)
    txt(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7), title, 34, tc, True)
    if sub:
        txt(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.4), sub, 15, sc, italic=True)


def table(slide, l, t, w, h, rows, cols, data, col_w=None, fsz=11, hdr_color=OCEAN):
    ts = slide.shapes.add_table(rows, cols, l, t, w, h)
    tbl = ts.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Inches(cw)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ''
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(fsz)
                p.font.name = 'Calibri'
                if r == 0:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                else:
                    p.font.color.rgb = CHARCOAL
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return ts


def pnum(slide, n, total, dark=False):
    c = RGBColor(0x88, 0xAA, 0xBB) if dark else MUTED
    txt(slide, Inches(12), Inches(7.0), Inches(1), Inches(0.3), f'{n}/{total}', 10, c, align=PP_ALIGN.RIGHT)


TOTAL = 13

# ================================================================
# SLIDE 1: TITLE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, CODE_BG)
rect(s, Inches(4.5), Inches(2.4), Inches(4.3), Inches(0.04), CORAL)
txt(s, Inches(1), Inches(1.0), Inches(11.3), Inches(1), 'ILOT COCO BEACH', 52, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(2.7), Inches(11.3), Inches(0.7), 'Technical Specification & Architecture', 24, CORAL_LIGHT, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5), 'Internal Document - Development Reference', 16, MUTED, align=PP_ALIGN.CENTER)

# Tech badges
techs = ['Vue.js 3', 'Node.js', 'PostgreSQL', 'TailwindCSS', 'GSAP', 'Prisma', 'JWT', 'Cloudflare R2']
for i, t in enumerate(techs):
    x = Inches(1.5) + Inches(i * 1.35)
    r = rect(s, x, Inches(5.0), Inches(1.15), Inches(0.4), RGBColor(0x2D, 0x3B, 0x4E), radius=0.15)
    txt(s, x + Inches(0.05), Inches(5.03), Inches(1.05), Inches(0.35), t, 11, OCEAN, True, PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(6.2), Inches(11.3), Inches(0.4), 'March 2026 | ~200h | 40 User Stories | 10 DB Tables', 12, MUTED, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 2: TECH STACK
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'Tech Stack', 'All tools and libraries')
pnum(s, 2, TOTAL)

categories = [
    ('FRONTEND', [
        ('Vue.js 3', 'Composition API + script setup'),
        ('Vite 8', 'Build tool + HMR'),
        ('TailwindCSS v4', '@theme tokens, RTL support'),
        ('Vue Router 4', 'SPA routing + hash anchors'),
        ('Vue I18n 11', 'FR / AR / EN translations'),
        ('Pinia', 'State management'),
        ('GSAP + ScrollTrigger', 'Scroll animations, entrance fx'),
        ('Lenis', 'Smooth scroll engine'),
        ('Swiper', 'Image carousels'),
    ], OCEAN),
    ('BACKEND', [
        ('Node.js + Fastify', 'REST API framework'),
        ('Prisma ORM', 'Type-safe DB queries, migrations'),
        ('JWT (jsonwebtoken)', 'Auth tokens (access + refresh)'),
        ('bcrypt', 'Password hashing'),
        ('node-cron', 'Scheduled jobs (flash sales)'),
        ('Multer + S3', 'File upload to object storage'),
        ('Passport.js', 'Google / Facebook OAuth'),
        ('Zod', 'Request validation'),
    ], CORAL),
    ('INFRA', [
        ('PostgreSQL 16', 'Primary database'),
        ('Cloudflare R2', 'Image/video storage (S3 API)'),
        ('Vercel', 'Frontend hosting (free)'),
        ('Railway / Hetzner', 'Backend hosting ($5-7/mo)'),
        ('Cloudflare', 'CDN + SSL + DDoS protection'),
        ('GitHub Actions', 'CI/CD pipeline'),
        ('Docker', 'Containerized deployment'),
    ], RGBColor(0x06, 0x96, 0x68)),
]

for ci, (cat_title, items, color) in enumerate(categories):
    x = Inches(0.5) + Inches(ci * 4.2)
    rect(s, x, Inches(1.7), Inches(3.9), Inches(5.3), WHITE, radius=0.06)
    rect(s, x, Inches(1.7), Inches(3.9), Inches(0.5), color, radius=0.06)
    rect(s, x, Inches(1.95), Inches(3.9), Inches(0.3), color)
    txt(s, x + Inches(0.2), Inches(1.75), Inches(3.5), Inches(0.45), cat_title, 16, WHITE, True, PP_ALIGN.CENTER)
    for i, (name, desc) in enumerate(items):
        y = Inches(2.35) + Inches(i * 0.5)
        txt(s, x + Inches(0.25), y, Inches(1.6), Inches(0.25), name, 11, CHARCOAL, True)
        txt(s, x + Inches(1.85), y, Inches(1.8), Inches(0.25), desc, 9, MUTED)


# ================================================================
# SLIDE 3: FRAMEWORK COMPARISON
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'Backend Framework Decision', 'Why Node.js over .NET, Django, Laravel')
pnum(s, 3, TOTAL)

data = [
    ['Criteria', 'Node.js (Fastify)', 'Django (Python)', '.NET Core (C#)', 'Laravel (PHP)'],
    ['Dev Speed (CRUD)', '★★★★★', '★★★★★', '★★★☆☆', '★★★★☆'],
    ['Hosting Cost', '★★★★★ ($5/mo)', '★★★★☆ ($5/mo)', '★★★☆☆ ($13+/mo)', '★★★★★ ($3/mo)'],
    ['Same Lang as FE', '★★★★★ (JS/TS)', '★★☆☆☆ (Python)', '★★☆☆☆ (C#)', '★★☆☆☆ (PHP)'],
    ['ORM Quality', 'Prisma (excellent)', 'Django ORM (great)', 'EF Core (great)', 'Eloquent (great)'],
    ['Auth Ecosystem', 'Passport.js + JWT', 'Django AllAuth', 'Identity + Azure', 'Sanctum + Breeze'],
    ['Cron Jobs', 'node-cron / Bull', 'Celery + Beat', 'Hangfire / Quartz', 'Built-in scheduler'],
    ['OVH Migration', '★★★★☆ (Docker)', '★★★★☆ (Docker)', '★★★☆☆ (.NET runtime)', '★★★★★ (native PHP)'],
    ['SCORE', '34/35', '29/35', '22/35', '29/35'],
]
table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.0), 9, 5, data,
      col_w=[2.3, 2.5, 2.5, 2.5, 2.5], fsz=11)

# Verdict
rect(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.3), CODE_BG, radius=0.06)
txt(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.4), 'DECISION: Node.js + Fastify + Prisma + PostgreSQL', 20, CORAL, True)
txt(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6),
    'One language everywhere (JS/TS). Cheapest to host. Fastest CRUD dev. Single dev can own full stack. '
    'If OVH shared hosting becomes mandatory, fallback to Laravel with same Vue.js frontend.',
    12, RGBColor(0xAA, 0xBB, 0xCC))


# ================================================================
# SLIDE 4: CLOUD HOSTING
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'Hosting & Infrastructure Costs', 'Phase 1 (managed) -> Phase 2 (VPS)')
pnum(s, 4, TOTAL)

data_m = [
    ['Provider', 'Type', 'Backend', 'DB', 'Cost/mo'],
    ['Railway', 'PaaS', 'Included', 'PG included', '$5-10'],
    ['Render', 'PaaS', 'Starter $7', 'PG $7', '$7-14'],
    ['Vercel+Supabase', 'Serverless', 'Edge funcs', 'Supabase free', '$0'],
    ['Fly.io', 'Containers', '3 free VMs', 'PG 3GB free', '$0-5'],
]
table(s, Inches(0.5), Inches(1.7), Inches(6), Inches(2.5), 5, 5, data_m, col_w=[1.5, 0.8, 1.2, 1.2, 0.8], fsz=11)

data_v = [
    ['Provider', 'Spec', 'Price/mo'],
    ['Hetzner CX22', '2 vCPU, 4GB, 40GB, 20TB', 'EUR 3.79'],
    ['OVH VPS Starter', '1 vCPU, 2GB, 20GB, unlim', 'EUR 3.50'],
    ['Contabo VPS S', '4 vCPU, 8GB, 50GB, 32TB', 'EUR 4.99'],
    ['DigitalOcean', '1 vCPU, 1GB, 25GB, 1TB', '$6.00'],
]
table(s, Inches(7), Inches(1.7), Inches(5.8), Inches(2.5), 5, 3, data_v, col_w=[2, 2.5, 1.0], fsz=11)

# Recommended stack
rect(s, Inches(0.5), Inches(4.5), Inches(6), Inches(2.7), WHITE, radius=0.06)
txt(s, Inches(0.8), Inches(4.6), Inches(5.4), Inches(0.4), 'Phase 1: Launch ($5-7/mo)', 16, OCEAN, True)
multiline(s, Inches(0.8), Inches(5.1), Inches(5.4), Inches(1.8), [
    'Frontend: Vercel (free, auto-deploy from GitHub)',
    'Backend: Railway Hobby ($5/mo, includes PostgreSQL)',
    'Storage: Cloudflare R2 (10GB free, S3-compatible)',
    'CDN/SSL: Cloudflare (free)',
    'Domain: ~$12/year',
    'TOTAL: ~$5-7/month',
], 12)

rect(s, Inches(7), Inches(4.5), Inches(5.8), Inches(2.7), WHITE, radius=0.06)
txt(s, Inches(7.3), Inches(4.6), Inches(5.2), Inches(0.4), 'Phase 2: OVH Migration (~EUR 4-7/mo)', 16, CORAL, True)
multiline(s, Inches(7.3), Inches(5.1), Inches(5.2), Inches(1.8), [
    'Hetzner CX22 or OVH VPS Comfort',
    'Docker: Node.js API + Nginx reverse proxy',
    'PostgreSQL self-hosted on same VPS',
    'R2 storage remains (or self-hosted MinIO)',
    'Cloudflare CDN remains (free)',
    'TOTAL: ~EUR 4-7/month all-in',
], 12)

# ================================================================
# SLIDE 5: DATABASE SCHEMA
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'Database Schema', '10 tables, PostgreSQL 16')
pnum(s, 5, TOTAL)

tables_data = [
    ('users', 'id (UUID PK), email, phone, password_hash,\nname, role (user|admin), provider\n(email|google|fb), created_at', OCEAN),
    ('menu_categories', 'id (UUID PK), name_fr, name_en,\nname_ar, sort_order', DRIFTWOOD),
    ('menu_items', 'id, category_id (FK), name_*, desc_*,\nprice_standard, price_extra, image_url,\nis_available, sort_order', DRIFTWOOD),
    ('emplacements', 'id, name_*, description_*, price,\nimage_url, capacity, sort_order', RGBColor(0x06, 0x96, 0x68)),
    ('images', 'id, url, section (hero|gallery|menu|\nemplacement), alt_text, sort_order', PURPLE),
    ('vouchers', 'id, code (unique), discount_percent,\ndiscount_amount, valid_from, valid_until,\nmax_uses, current_uses, is_active', CORAL),
    ('flash_sales', 'id, title_*, description_*, discount_%,\nstarts_at, ends_at, is_active, image_url', RED),
    ('reviews', 'id, user_id (FK), rating (1-5), comment,\nis_approved, is_visible, created_at', YELLOW),
    ('reservations', 'id, user_id (FK nullable), name, phone,\ndate, guests, message, voucher_id (FK),\nstatus (pending|confirmed|cancelled)', OCEAN),
    ('site_settings', 'key (PK), value (TEXT)\nStores: phone, address, lat/lng,\nabout_text, seo_keywords, etc.', MUTED),
]

for i, (name, cols, color) in enumerate(tables_data):
    col = i % 5
    row = i // 5
    x = Inches(0.3) + Inches(col * 2.6)
    y = Inches(1.7) + Inches(row * 2.9)
    r = rect(s, x, y, Inches(2.4), Inches(2.6), WHITE, radius=0.06)
    rect(s, x, y, Inches(2.4), Inches(0.45), color, radius=0.06)
    rect(s, x, y + Inches(0.22), Inches(2.4), Inches(0.25), color)
    txt(s, x + Inches(0.1), y + Inches(0.07), Inches(2.2), Inches(0.35), name, 13, WHITE, True, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), y + Inches(0.55), Inches(2.2), Inches(2.0), cols, 9, MUTED)


# ================================================================
# SLIDE 6: API ENDPOINTS
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'REST API Endpoints', 'Fastify routes grouped by resource')
pnum(s, 6, TOTAL)

endpoints = [
    ['Group', 'Method', 'Endpoint', 'Auth', 'Description'],
    ['Auth', 'POST', '/api/auth/register', 'Public', 'Sign up (email/phone+pass)'],
    ['Auth', 'POST', '/api/auth/login', 'Public', 'Login, returns JWT'],
    ['Auth', 'POST', '/api/auth/google', 'Public', 'Google OAuth callback'],
    ['Auth', 'POST', '/api/auth/facebook', 'Public', 'Facebook OAuth callback'],
    ['Menu', 'GET', '/api/menu', 'Public', 'Get all categories + items'],
    ['Menu', 'POST', '/api/admin/menu/items', 'Admin', 'Create menu item'],
    ['Menu', 'PUT', '/api/admin/menu/items/:id', 'Admin', 'Update menu item'],
    ['Menu', 'DELETE', '/api/admin/menu/items/:id', 'Admin', 'Delete menu item'],
    ['Images', 'GET', '/api/images/:section', 'Public', 'Get images by section'],
    ['Images', 'POST', '/api/admin/images', 'Admin', 'Upload image'],
    ['Images', 'DELETE', '/api/admin/images/:id', 'Admin', 'Delete image'],
    ['Vouchers', 'GET', '/api/vouchers/active', 'Public', 'Get active vouchers'],
    ['Vouchers', 'POST', '/api/admin/vouchers', 'Admin', 'Create voucher'],
    ['Flash', 'GET', '/api/flash-sales/active', 'Public', 'Get active flash sales'],
    ['Flash', 'POST', '/api/admin/flash-sales', 'Admin', 'Create flash sale'],
    ['Reviews', 'GET', '/api/reviews', 'Public', 'Get approved reviews'],
    ['Reviews', 'POST', '/api/reviews', 'User', 'Submit a review'],
    ['Reviews', 'PUT', '/api/admin/reviews/:id', 'Admin', 'Approve/reject review'],
    ['Reservations', 'POST', '/api/reservations', 'Public', 'Submit reservation'],
    ['Settings', 'GET', '/api/settings', 'Public', 'Get site settings'],
    ['Settings', 'PUT', '/api/admin/settings', 'Admin', 'Update site settings'],
]
table(s, Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.7), 22, 5, endpoints,
      col_w=[1.0, 0.8, 3.5, 0.8, 3.0], fsz=10)


# ================================================================
# SLIDE 7: USER STORIES FULL
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'User Stories Backlog', '40 stories across 8 epics')
pnum(s, 7, TOTAL)

stories = [
    ['Epic', '#', 'Story', 'Priority', 'Hours'],
    ['Public', '12', 'Hero video, about, gallery, menu, emplacements, vouchers, flash sales, reviews, social, map, i18n, responsive', 'HIGH', '50h'],
    ['Auth', '4', 'Register (email/phone), Login, Google/FB OAuth, Password reset', 'HIGH', '17h'],
    ['Reservation', '3', 'Form + WhatsApp send, Apply voucher, View history', 'HIGH', '11h'],
    ['Reviews', '2', 'Submit review (logged in), Display approved only', 'HIGH', '6h'],
    ['Admin Content', '6', 'Login, Menu CRUD, Image mgmt, Hero video, Emplacements, Settings', 'HIGH', '30h'],
    ['Admin Promos', '4', 'Voucher CRUD, Flash sale scheduling, Auto-activate/deactivate', 'HIGH', '18h'],
    ['Admin Moderation', '4', 'View/approve/reject reviews, View reservations, View users', 'MED', '11h'],
    ['Infra', '5', 'DB schema + migrations, API + auth middleware, File storage, Cron jobs, CI/CD', 'HIGH', '26h'],
    ['QA + Deploy', '-', 'Testing, bug fixes, responsive polish, production deploy, client training', 'HIGH', '30h'],
    ['TOTAL', '40', '', '', '199h'],
]
table(s, Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.2), 11, 5, stories,
      col_w=[1.8, 0.5, 7.0, 0.8, 0.8], fsz=12)


# ================================================================
# SLIDE 8: CRON JOBS & BACKGROUND TASKS
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'Cron Jobs & Automation', 'Scheduled tasks for flash sales and vouchers')
pnum(s, 8, TOTAL)

# Flash sale cron
rect(s, Inches(0.5), Inches(1.8), Inches(6), Inches(5.0), WHITE, radius=0.06)
rect(s, Inches(0.5), Inches(1.8), Inches(6), Inches(0.5), CORAL, radius=0.06)
rect(s, Inches(0.5), Inches(2.05), Inches(6), Inches(0.28), CORAL)
txt(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(0.4), 'Flash Sale Cron', 16, WHITE, True)
multiline(s, Inches(0.8), Inches(2.5), Inches(5.4), Inches(4.0), [
    'Schedule: Every minute (* * * * *)',
    '',
    'Logic:',
    '1. Query flash_sales WHERE starts_at <= NOW() AND is_active = false',
    '   -> SET is_active = true (activate)',
    '',
    '2. Query flash_sales WHERE ends_at <= NOW() AND is_active = true',
    '   -> SET is_active = false (deactivate)',
    '',
    'Implementation: node-cron in Fastify plugin',
    'Alternative: Railway/Render cron if available',
    '',
    'Frontend: Countdown timer reads ends_at,',
    'calculates remaining time client-side (no API polling)',
], 11, CHARCOAL)

# Voucher cron
rect(s, Inches(7), Inches(1.8), Inches(5.8), Inches(5.0), WHITE, radius=0.06)
rect(s, Inches(7), Inches(1.8), Inches(5.8), Inches(0.5), GREEN, radius=0.06)
rect(s, Inches(7), Inches(2.05), Inches(5.8), Inches(0.28), GREEN)
txt(s, Inches(7.2), Inches(1.85), Inches(5.4), Inches(0.4), 'Voucher Expiry Cron', 16, WHITE, True)
multiline(s, Inches(7.3), Inches(2.5), Inches(5.2), Inches(4.0), [
    'Schedule: Every hour (0 * * * *)',
    '',
    'Logic:',
    '1. Query vouchers WHERE valid_until <= NOW()',
    '   AND is_active = true',
    '   -> SET is_active = false',
    '',
    '2. Query vouchers WHERE current_uses >= max_uses',
    '   AND is_active = true',
    '   -> SET is_active = false',
    '',
    'Validation also happens at reservation time:',
    'API checks voucher validity before accepting',
    '',
    'Admin can manually activate/deactivate anytime',
], 11, CHARCOAL)


# ================================================================
# SLIDE 9: AUTH FLOW
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'Authentication Architecture', 'JWT + OAuth2 flow')
pnum(s, 9, TOTAL)

flows = [
    ('Email/Phone Registration', [
        'Client sends POST /auth/register {email, password, name}',
        'Server validates with Zod, hashes password (bcrypt, 12 rounds)',
        'Creates user record, returns JWT access + refresh tokens',
        'Access token: 15min expiry, stored in memory',
        'Refresh token: 7d expiry, httpOnly cookie',
    ], OCEAN),
    ('Google/Facebook OAuth', [
        'Client redirects to /auth/google (or /facebook)',
        'Provider returns auth code to callback URL',
        'Server exchanges code for provider access token',
        'Fetches user profile from provider API',
        'Creates/links user, returns JWT tokens',
    ], CORAL),
    ('Admin Auth', [
        'Same JWT flow but with role check middleware',
        'admin routes check: req.user.role === "admin"',
        'Admin accounts created via DB seed or API',
        'No self-registration for admin (security)',
        'Consider: 2FA via TOTP for admin accounts',
    ], PURPLE),
]

for i, (title, steps, color) in enumerate(flows):
    x = Inches(0.5) + Inches(i * 4.2)
    rect(s, x, Inches(1.7), Inches(3.9), Inches(5.3), WHITE, radius=0.06)
    rect(s, x, Inches(1.7), Inches(3.9), Inches(0.5), color, radius=0.06)
    rect(s, x, Inches(1.95), Inches(3.9), Inches(0.28), color)
    txt(s, x + Inches(0.2), Inches(1.75), Inches(3.5), Inches(0.4), title, 14, WHITE, True)
    for j, step in enumerate(steps):
        y = Inches(2.4) + Inches(j * 0.55)
        c = circle(s, x + Inches(0.2), y + Inches(0.04), Inches(0.22), color)
        tf = c.text_frame
        tf.paragraphs[0].text = str(j + 1)
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        txt(s, x + Inches(0.5), y, Inches(3.2), Inches(0.5), step, 10, CHARCOAL)


# ================================================================
# SLIDE 10: FRONTEND ARCHITECTURE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'Frontend Architecture', 'Vue 3 component structure + state management')
pnum(s, 10, TOTAL)

rect(s, Inches(0.5), Inches(1.7), Inches(6), Inches(5.3), WHITE, radius=0.06)
txt(s, Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.4), 'Component Tree', 16, OCEAN, True)
tree = [
    'App.vue (Lenis smooth scroll, i18n dir/lang)',
    '  HomeView.vue',
    '    NavBar.vue (IntersectionObserver, lang switcher)',
    '    HeroSection.vue (Ken Burns, GSAP stagger)',
    '    AboutSection.vue (Parallax, feature cards)',
    '    GallerySection.vue (Masonry, lightbox, clip-path)',
    '    VideoSection.vue (Custom player, scroll trigger)',
    '    ExperienceSection.vue (Horiz scroll, 3D tilt)',
    '    ReservationSection.vue (WhatsApp form builder)',
    '    LocationSection.vue (Maps embed, contact cards)',
    '    FooterSection.vue (Wave SVG, social links)',
    '  AdminLogin.vue (Supabase/JWT auth)',
    '  AdminDashboard.vue (Image CRUD, menu mgmt)',
]
multiline(s, Inches(0.7), Inches(2.3), Inches(5.6), Inches(4.5), tree, 11, CHARCOAL)

rect(s, Inches(7), Inches(1.7), Inches(5.8), Inches(2.4), WHITE, radius=0.06)
txt(s, Inches(7.2), Inches(1.8), Inches(5.4), Inches(0.4), 'State Management (Pinia)', 16, OCEAN, True)
multiline(s, Inches(7.2), Inches(2.3), Inches(5.4), Inches(1.6), [
    'useAuthStore - user session, JWT tokens, role',
    'useMenuStore - categories, items, fetch/cache',
    'useGalleryStore - images by section, admin CRUD',
    'useSettingsStore - site config, SEO, contact info',
    'usePromoStore - vouchers, flash sales, timers',
], 11, CHARCOAL)

rect(s, Inches(7), Inches(4.3), Inches(5.8), Inches(2.7), WHITE, radius=0.06)
txt(s, Inches(7.2), Inches(4.4), Inches(5.4), Inches(0.4), 'Key Composables', 16, OCEAN, True)
multiline(s, Inches(7.2), Inches(4.9), Inches(5.4), Inches(1.8), [
    'useConfig() - env vars (phone, address, socials)',
    'useScrollAnimation() - GSAP ScrollTrigger setup',
    'useWhatsApp() - builds + encodes WA message URL',
    'useCountdown(endDate) - reactive timer for flash sales',
    'useImageUpload() - S3/R2 upload with progress',
    'useApi() - fetch wrapper with JWT auto-refresh',
], 11, CHARCOAL)


# ================================================================
# SLIDE 11: TIMELINE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'Sprint Plan', '10 weeks, 2-week sprints')
pnum(s, 11, TOTAL)

sprints = [
    ('Sprint 1\nWeek 1-2', 'Infrastructure + Public Site', [
        'DB schema + Prisma migrations',
        'Fastify API scaffold + auth middleware',
        'Hero, About, Gallery, Video sections',
        'Lenis + GSAP scroll animations',
        'i18n setup (FR/AR/EN)',
    ], OCEAN),
    ('Sprint 2\nWeek 3-4', 'Auth + Dynamic Content', [
        'User registration (email/phone)',
        'Google + Facebook OAuth',
        'Menu page (Standard/Extra)',
        'Emplacements page',
        'Review display + submission',
    ], RGBColor(0x06, 0x96, 0x68)),
    ('Sprint 3\nWeek 5-6', 'Backoffice Core', [
        'Admin dashboard layout',
        'Menu CRUD (categories + items)',
        'Image manager (upload/reorder/delete)',
        'Emplacement editor',
        'Site settings editor',
    ], DRIFTWOOD),
    ('Sprint 4\nWeek 7-8', 'Promos + Automation', [
        'Voucher CRUD + validation',
        'Flash sale scheduling + cron',
        'Countdown timer component',
        'Review moderation panel',
        'Reservation viewer',
    ], CORAL),
    ('Sprint 5\nWeek 9-10', 'QA + Launch', [
        'Responsive testing (all breakpoints)',
        'Cross-browser testing',
        'Performance optimization (lazy load)',
        'Production deployment',
        'Client training session',
    ], PURPLE),
]

for i, (period, title, tasks, color) in enumerate(sprints):
    x = Inches(0.3) + Inches(i * 2.6)
    rect(s, x, Inches(1.7), Inches(2.4), Inches(5.3), WHITE, radius=0.06)
    rect(s, x, Inches(1.7), Inches(2.4), Inches(0.9), color, radius=0.06)
    rect(s, x, Inches(2.35), Inches(2.4), Inches(0.28), color)
    txt(s, x + Inches(0.1), Inches(1.75), Inches(2.2), Inches(0.35), period, 11, WHITE, True, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.15), Inches(2.2), Inches(0.3), title, 10, RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)
    for j, task in enumerate(tasks):
        y = Inches(2.8) + Inches(j * 0.45)
        c = circle(s, x + Inches(0.15), y + Inches(0.06), Inches(0.15), color)
        txt(s, x + Inches(0.4), y, Inches(1.85), Inches(0.4), task, 9, CHARCOAL)


# ================================================================
# SLIDE 12: COST SUMMARY
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, SAND)
header(s, 'Cost Summary', 'Development + monthly hosting')
pnum(s, 12, TOTAL)

data_dev = [
    ['Approach', 'Rate', 'Hours', 'Total'],
    ['Freelance (Tunisia)', '30-50 TND/h', '200h', '6,000-10,000 TND'],
    ['Freelance (International)', '$25-50/h', '200h', '$5,000-10,000'],
    ['Small Agency', '$40-80/h', '200h', '$8,000-16,000'],
]
txt(s, Inches(0.5), Inches(1.6), Inches(6), Inches(0.3), 'Development Cost', 16, OCEAN, True)
table(s, Inches(0.5), Inches(2.0), Inches(6), Inches(1.8), 4, 4, data_dev, col_w=[2.2, 1.3, 0.7, 2.2], fsz=12)

data_host = [
    ['Service', 'Phase 1', 'Phase 2 (OVH)'],
    ['Backend hosting', 'Railway $5/mo', 'Hetzner EUR 3.79/mo'],
    ['Frontend hosting', 'Vercel free', 'Same VPS'],
    ['Database', 'Railway PG (included)', 'Self-hosted PG'],
    ['File storage', 'Cloudflare R2 (free)', 'R2 or MinIO'],
    ['CDN + SSL', 'Cloudflare (free)', 'Cloudflare (free)'],
    ['Domain', '~$12/year', '~$12/year'],
    ['TOTAL', '$5-7/month', 'EUR 4-7/month'],
]
txt(s, Inches(7), Inches(1.6), Inches(6), Inches(0.3), 'Monthly Hosting Cost', 16, OCEAN, True)
table(s, Inches(7), Inches(2.0), Inches(5.8), Inches(3.5), 8, 3, data_host, col_w=[1.8, 2.0, 2.0], fsz=11)

# Client pricing suggestion
rect(s, Inches(0.5), Inches(4.3), Inches(6), Inches(2.8), CODE_BG, radius=0.06)
txt(s, Inches(0.8), Inches(4.4), Inches(5.4), Inches(0.4), 'Suggested Client Pricing', 16, CORAL, True)
multiline(s, Inches(0.8), Inches(4.9), Inches(5.4), Inches(1.8), [
    'Option A: 10,000 TND flat + 200 TND/mo maintenance',
    'Option B: Phase 1 MVP 6,000 TND + Phase 2 4,000 TND',
    '',
    'Includes: Source code ownership, 3 months support,',
    'backoffice training, documentation, first year hosting.',
    '',
    'Margin: ~40-60% depending on your hourly rate.',
], 12, RGBColor(0xAA, 0xBB, 0xCC))


# ================================================================
# SLIDE 13: END
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, CODE_BG)
rect(s, Inches(4.8), Inches(2.2), Inches(3.7), Inches(0.04), CORAL)
txt(s, Inches(1), Inches(1.0), Inches(11.3), Inches(0.8), 'ILOT COCO BEACH', 48, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(2.5), Inches(11.3), Inches(0.5), 'Technical Reference Document', 20, CORAL_LIGHT, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.5), Inches(11.3), Inches(0.5), 'Ready to build.', 18, MUTED, align=PP_ALIGN.CENTER, italic=True)

output = r'C:\Personal\ilot-cocobeach\Ilot_CocoBeach_TECHNICAL.pptx'
prs.save(output)
print(f'Technical presentation saved: {output}')
print(f'Slides: {len(prs.slides)}')
