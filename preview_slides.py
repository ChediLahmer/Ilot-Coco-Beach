import subprocess, sys
# Convert a few slides to images using LibreOffice if available, else just verify structure
from pptx import Presentation

for name in ['Ilot_CocoBeach_CLIENT.pptx', 'Ilot_CocoBeach_TECHNICAL.pptx']:
    path = f'C:\\Personal\\ilot-cocobeach\\{name}'
    prs = Presentation(path)
    print(f'\n=== {name} ({len(prs.slides)} slides) ===')
    for i, slide in enumerate(prs.slides):
        shape_count = len(slide.shapes)
        texts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        texts.append(t)
        title = texts[0] if texts else '(empty)'
        print(f'  Slide {i+1}: {title[:60]:<60} | {shape_count} shapes')
